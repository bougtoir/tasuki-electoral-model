#!/usr/bin/env python3
"""Generate English and Japanese PPTX files with figures.
- Code-generated plots (Fig 2-5): embedded as images
- Flow/conceptual diagrams (Fig 1, Fig 6): editable PowerPoint shapes
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

FIGS = '/home/ubuntu/figures'

# ── Color palette ──
C_CANDIDATE = RGBColor(0x21, 0x96, 0xF3)
C_ELECTION  = RGBColor(0x4C, 0xAF, 0x50)
C_TERM      = RGBColor(0xFF, 0x98, 0x00)
C_EVAL      = RGBColor(0xE9, 0x1E, 0x63)
C_WEIGHT    = RGBColor(0x9C, 0x27, 0xB0)
C_VOTER     = RGBColor(0x60, 0x7D, 0x8B)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY      = RGBColor(0x99, 0x99, 0x99)
C_BROWN     = RGBColor(0x79, 0x55, 0x48)
C_CYAN      = RGBColor(0x00, 0xBC, 0xD4)
C_TASUKI       = RGBColor(0xE9, 0x1E, 0x63)

def add_rounded_box(slide, left, top, width, height, text, fill_color, font_size=10):
    """Add a rounded rectangle with centered white text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = C_WHITE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
    tf.paragraphs[0].space_before = Pt(2)
    return shape

def add_arrow_shape(slide, left, top, width, height, color=C_DARK, direction='right'):
    """Add an arrow shape for better visibility than connectors."""
    shape_map = {
        'right': MSO_SHAPE.RIGHT_ARROW,
        'left':  MSO_SHAPE.LEFT_ARROW,
        'down':  MSO_SHAPE.DOWN_ARROW,
        'up':    MSO_SHAPE.UP_ARROW,
    }
    shape = slide.shapes.add_shape(shape_map[direction], left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=11,
                 bold=False, italic=False, color=C_DARK, alignment=PP_ALIGN.LEFT):
    """Add a text box with explicit no-fill background."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_circle(slide, left, top, diameter, text, fill_color, font_size=8):
    """Add a circle shape with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = C_WHITE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
    return shape

def build_slide1_en(prs):
    """Fig 1: Conceptual Overview - editable (English)
    Z-order: arrows (back) -> boxes (front) -> text labels (top)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bw, bh = Inches(2.6), Inches(0.95)
    y_top = Inches(1.2)
    y_bot = Inches(3.0)
    x1 = Inches(0.4)
    x2 = Inches(3.5)
    x3 = Inches(6.6)

    # LAYER 1: Arrows FIRST (behind the boxes)
    arrow_y_top = y_top + bh // 2 - Inches(0.1)
    arrow_y_bot = y_bot + bh // 2 - Inches(0.1)
    arrow_h = Inches(0.2)
    arrow_gap = Inches(0.05)
    add_arrow_shape(slide, x1 + bw + arrow_gap, arrow_y_top,
                    x2 - x1 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'right')
    add_arrow_shape(slide, x2 + bw + arrow_gap, arrow_y_top,
                    x3 - x2 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'right')
    add_arrow_shape(slide, x3 + bw // 2 - Inches(0.1), y_top + bh + arrow_gap,
                    Inches(0.2), y_bot - y_top - bh - 2 * arrow_gap, C_DARK, 'down')
    add_arrow_shape(slide, x2 + bw + arrow_gap, arrow_y_bot,
                    x3 - x2 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'left')
    add_arrow_shape(slide, x1 + bw + arrow_gap, arrow_y_bot,
                    x2 - x1 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'left')
    add_arrow_shape(slide, x1 + bw // 2 - Inches(0.12), y_top + bh + arrow_gap,
                    Inches(0.24), y_bot - y_top - bh - 2 * arrow_gap, C_TASUKI, 'up')

    # LAYER 2: Phase boxes ON TOP of arrows
    add_rounded_box(slide, x1, y_top, bw, bh,
                    'Phase 1\nPledge Declaration\nw₁, w₂, …, wₖ (Σ=1)',
                    C_CANDIDATE, 10)
    add_rounded_box(slide, x2, y_top, bw, bh,
                    'Phase 2\nElection\n(weighted ballots)',
                    C_ELECTION, 10)
    add_rounded_box(slide, x3, y_top, bw, bh,
                    'Phase 3\nTerm in Office\n(interim eval: O/P/E)',
                    C_TERM, 10)
    add_rounded_box(slide, x1, y_bot, bw, bh,
                    'Phase 6\nTrust Update\nτ(S) → next election',
                    C_VOTER, 10)
    add_rounded_box(slide, x2, y_bot, bw, bh,
                    'Phase 5\nAccountability Score\nS = Σ wⱼ·fⱼ',
                    C_WEIGHT, 10)
    add_rounded_box(slide, x3, y_bot, bw, bh,
                    'Phase 4\nFulfillment Eval\nfⱼ ∈ [0,1] per pledge',
                    C_EVAL, 10)

    # LAYER 3: Title and text labels (on top of everything)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 'Figure 1. TASUKI Conceptual Overview', font_size=20,
                 bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.3),
                 'Trust-Adjusted Scoring with Unified Knowledge Integration (襟) — 6-Phase Cycle',
                 font_size=12, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(2.1), Inches(0.7), Inches(0.4),
                 'Next\nCycle', font_size=9, bold=True, color=C_TASUKI,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.35),
                 'Core Equations', font_size=14, bold=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(0.3),
                 'Accountability Score:  Sᶜ = Σ wⱼ · fⱼ  where Σwⱼ = 1',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.05), Inches(9), Inches(0.3),
                 'Candidate Trust:  τᶜ⁽ᵗ⁺¹⁾ = ω(Sᶜ⁽ᵗ⁾)  where ω: [0,1] → [τₘᵢₙ, τₘₐₓ]',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.4), Inches(9), Inches(0.3),
                 'Effective Vote:  V_eff(c) = τᶜ · nᶜ  (raw votes for candidate c)',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.5),
                 'Figure 1. Conceptual overview of the TASUKI mechanism showing the six phases of each electoral cycle: '
                 'pledge declaration, election, term in office, fulfillment evaluation, accountability score computation, '
                 'and trust coefficient update.',
                 font_size=9, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)

def build_slide1_ja(prs):
    """Fig 1: Conceptual Overview - editable (Japanese)
    Z-order: arrows (back) -> boxes (front) -> text labels (top)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bw, bh = Inches(2.6), Inches(0.95)
    y_top = Inches(1.2)
    y_bot = Inches(3.0)
    x1 = Inches(0.4)
    x2 = Inches(3.5)
    x3 = Inches(6.6)

    # LAYER 1: Arrows FIRST
    arrow_y_top = y_top + bh // 2 - Inches(0.1)
    arrow_y_bot = y_bot + bh // 2 - Inches(0.1)
    arrow_h = Inches(0.2)
    arrow_gap = Inches(0.05)
    add_arrow_shape(slide, x1 + bw + arrow_gap, arrow_y_top,
                    x2 - x1 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'right')
    add_arrow_shape(slide, x2 + bw + arrow_gap, arrow_y_top,
                    x3 - x2 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'right')
    add_arrow_shape(slide, x3 + bw // 2 - Inches(0.1), y_top + bh + arrow_gap,
                    Inches(0.2), y_bot - y_top - bh - 2 * arrow_gap, C_DARK, 'down')
    add_arrow_shape(slide, x2 + bw + arrow_gap, arrow_y_bot,
                    x3 - x2 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'left')
    add_arrow_shape(slide, x1 + bw + arrow_gap, arrow_y_bot,
                    x2 - x1 - bw - 2 * arrow_gap, arrow_h, C_DARK, 'left')
    add_arrow_shape(slide, x1 + bw // 2 - Inches(0.12), y_top + bh + arrow_gap,
                    Inches(0.24), y_bot - y_top - bh - 2 * arrow_gap, C_TASUKI, 'up')

    # LAYER 2: Phase boxes
    add_rounded_box(slide, x1, y_top, bw, bh,
                    'フェーズ1\n公約宣言\nw₁, w₂, …, wₖ (Σ=1)',
                    C_CANDIDATE, 10)
    add_rounded_box(slide, x2, y_top, bw, bh,
                    'フェーズ2\n選挙\n（加重投票）',
                    C_ELECTION, 10)
    add_rounded_box(slide, x3, y_top, bw, bh,
                    'フェーズ3\n任期\n（中間評価: O/P/E）',
                    C_TERM, 10)
    add_rounded_box(slide, x1, y_bot, bw, bh,
                    'フェーズ6\n信任係数更新\nτ(S) → 次回選挙',
                    C_VOTER, 10)
    add_rounded_box(slide, x2, y_bot, bw, bh,
                    'フェーズ5\n説明責任得点\nS = Σ wⱼ·fⱼ',
                    C_WEIGHT, 10)
    add_rounded_box(slide, x3, y_bot, bw, bh,
                    'フェーズ4\n実現度評価\nfⱼ ∈ [0,1] 公約ごと',
                    C_EVAL, 10)

    # LAYER 3: Title and text labels
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 '図1. TASUKI 概念的概要', font_size=20,
                 bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.3),
                 '統合的知識に基づく信頼調整型評価システム（襟）— 6フェーズサイクル',
                 font_size=12, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(2.1), Inches(0.7), Inches(0.4),
                 '次の\nサイクル', font_size=9, bold=True, color=C_TASUKI,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.35),
                 '主要数式', font_size=14, bold=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(0.3),
                 '説明責任得点:  Sᶜ = Σ wⱼ · fⱼ  (Σwⱼ = 1)',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.05), Inches(9), Inches(0.3),
                 '信任係数:  τᶜ⁽ᵗ⁺¹⁾ = ω(Sᶜ⁽ᵗ⁾)  ω: [0,1] → [τₘᵢₙ, τₘₐₓ]',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.4), Inches(9), Inches(0.3),
                 '有効得票数:  V_eff(c) = τᶜ · nᶜ  (候補者cの生票数)',
                 font_size=12, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.5),
                 '図1. TASUKIメカニズムの概念的概要。各選挙サイクルの6フェーズを示す：'
                 '公約宣言、選挙、任期、実現度評価、説明責任得点算出、信任係数更新。',
                 font_size=9, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)

def build_image_slide(prs, img_path, title, caption):
    """Add a slide with an auto-scaled embedded image figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
                 title, font_size=18, bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)

    # Auto-scale image to fit available space
    max_img_w = Inches(9.2)
    max_img_h = Inches(5.7)  # leave room for title and caption
    img_top = Inches(0.7)

    # Get image dimensions for aspect ratio
    with Image.open(img_path) as im:
        img_w_px, img_h_px = im.size
    aspect = img_w_px / img_h_px

    # Scale to fit
    if max_img_w / aspect <= max_img_h:
        final_w = max_img_w
        final_h = int(max_img_w / aspect)
    else:
        final_h = max_img_h
        final_w = int(max_img_h * aspect)

    # Center horizontally
    slide_w = Inches(10)
    img_left = (slide_w - final_w) // 2

    slide.shapes.add_picture(img_path, img_left, img_top, width=final_w, height=final_h)

    # Caption below the image
    caption_y = img_top + final_h + Inches(0.1)
    if caption_y > Inches(6.8):
        caption_y = Inches(6.8)
    add_text_box(slide, Inches(0.3), caption_y, Inches(9.4), Inches(0.6),
                 caption, font_size=9, italic=True, color=C_GRAY,
                 alignment=PP_ALIGN.CENTER)

def build_slide6_en(prs):
    """Fig 6: Positioning Map - editable (English)
    Z-order: bg rect -> quadrant labels -> circles -> labels -> axis labels -> caption
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # LAYER 1: Plot background rectangle (added FIRST = behind everything)
    plot_left = Inches(1.0)
    plot_top = Inches(0.8)
    plot_w = Inches(8.5)
    plot_h = Inches(5.8)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, plot_left, plot_top, plot_w, plot_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    bg.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    bg.line.width = Pt(1)

    # LAYER 2: Quadrant labels (behind circles)
    add_text_box(slide, Inches(1.2), Inches(0.9), Inches(1.5), Inches(0.4),
                 'Traditional\nDemocracy', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(8.0), Inches(0.9), Inches(1.5), Inches(0.4),
                 'Radical\nReform', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(1.2), Inches(6.0), Inches(1.5), Inches(0.4),
                 'Accountability\nFocused', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(8.0), Inches(6.0), Inches(1.5), Inches(0.4),
                 'Mechanism\nDesign', font_size=9, italic=True, color=C_GRAY)

    def map_pos(x_norm, y_norm):
        px = plot_left + int(x_norm * plot_w)
        py = plot_top + int((0.75 - y_norm) / (0.75 - 0.1) * plot_h)
        return px, py

    # Positions adjusted to avoid overlaps
    approaches = [
        ('Standard\nElection',              0.08, 0.55, C_VOTER,     Inches(0.55)),
        ('Retrospective\nVoting (Fiorina)',  0.08, 0.28, C_BROWN,    Inches(0.50)),
        ('Barro-Ferejohn\nModel',            0.24, 0.38, C_ELECTION,  Inches(0.50)),
        ('Liquid\nDemocracy',                0.42, 0.55, C_CANDIDATE, Inches(0.55)),
        ('Weighted\nVoting',                 0.58, 0.55, C_CYAN,      Inches(0.50)),
        ('Quadratic\nVoting',                0.74, 0.55, C_WEIGHT,    Inches(0.55)),
        ('Futarchy',                         0.85, 0.65, C_TERM,      Inches(0.50)),
        ('TASUKI\n(This Paper)',             0.50, 0.20, C_TASUKI,       Inches(0.75)),
    ]

    # LAYER 3: ALL circles first (so labels appear on top)
    circle_positions = []
    for name, xn, yn, color, diam in approaches:
        px, py = map_pos(xn, yn)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, px - diam // 2, py - diam // 2, diam, diam)
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.color.rgb = C_WHITE
        circ.line.width = Pt(2)
        circle_positions.append((name, px, py, color, diam))

    # LAYER 4: ALL labels on top of circles
    for name, px, py, color, diam in circle_positions:
        fs = 11 if 'TASUKI' in name else 9
        bld = 'TASUKI' in name
        if 'TASUKI' in name:
            label_y = py - diam // 2 - Inches(0.35)
        else:
            label_y = py + diam // 2 + Inches(0.05)
        add_text_box(slide, px - Inches(0.8), label_y, Inches(1.6), Inches(0.4),
                     name, font_size=fs, bold=bld, color=color, alignment=PP_ALIGN.CENTER)

    # LAYER 5: Title and axis labels (on top of everything)
    add_text_box(slide, Inches(0.5), Inches(0.1), Inches(9), Inches(0.45),
                 'Figure 6. Theoretical Positioning of TASUKI', font_size=18,
                 bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.3),
                 'Degree of Departure from One-Person-One-Vote  →',
                 font_size=11, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(0.9), Inches(0.7), Inches(0.5),
                 '←\nRetro.', font_size=10, italic=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(5.7), Inches(0.7), Inches(0.5),
                 'Prosp.\n→', font_size=10, italic=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)

    # LAYER 6: Caption
    add_text_box(slide, Inches(0.3), Inches(7.05), Inches(9.4), Inches(0.4),
                 'Figure 6. Positioning of TASUKI relative to existing electoral reform approaches. '
                 'Horizontal axis: departure from OPOV. Vertical axis: temporal orientation (retrospective to prospective).',
                 font_size=9, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)

def build_slide6_ja(prs):
    """Fig 6: Positioning Map - editable (Japanese)
    Z-order: bg rect -> quadrant labels -> circles -> labels -> axis labels -> caption
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # LAYER 1: Plot background rectangle
    plot_left = Inches(1.0)
    plot_top = Inches(0.8)
    plot_w = Inches(8.5)
    plot_h = Inches(5.8)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, plot_left, plot_top, plot_w, plot_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    bg.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    bg.line.width = Pt(1)

    # LAYER 2: Quadrant labels
    add_text_box(slide, Inches(1.2), Inches(0.9), Inches(1.5), Inches(0.4),
                 '伝統的\n民主主義', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(8.0), Inches(0.9), Inches(1.5), Inches(0.4),
                 '急進的\n改革', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(1.2), Inches(6.0), Inches(1.5), Inches(0.4),
                 '説明責任\n重視', font_size=9, italic=True, color=C_GRAY)
    add_text_box(slide, Inches(8.0), Inches(6.0), Inches(1.5), Inches(0.4),
                 'メカニズム\nデザイン', font_size=9, italic=True, color=C_GRAY)

    def map_pos(x_norm, y_norm):
        px = plot_left + int(x_norm * plot_w)
        py = plot_top + int((0.75 - y_norm) / (0.75 - 0.1) * plot_h)
        return px, py

    # Positions adjusted to avoid overlaps (same as English)
    approaches_ja = [
        ('標準選挙',                        0.08, 0.55, C_VOTER,     Inches(0.55)),
        ('回顧的投票\n(Fiorina)',        0.08, 0.28, C_BROWN,     Inches(0.50)),
        ('Barro-Ferejohn\nモデル',              0.24, 0.38, C_ELECTION,  Inches(0.50)),
        ('Liquid\nDemocracy',                                0.42, 0.55, C_CANDIDATE, Inches(0.55)),
        ('加重投票',                         0.58, 0.55, C_CYAN,      Inches(0.50)),
        ('Quadratic\nVoting',                                0.74, 0.55, C_WEIGHT,    Inches(0.55)),
        ('Futarchy',                                         0.85, 0.65, C_TERM,      Inches(0.50)),
        ('TASUKI\n(本論文)',                     0.50, 0.20, C_TASUKI,       Inches(0.75)),
    ]

    # LAYER 3: ALL circles first
    circle_positions = []
    for name, xn, yn, color, diam in approaches_ja:
        px, py = map_pos(xn, yn)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, px - diam // 2, py - diam // 2, diam, diam)
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.color.rgb = C_WHITE
        circ.line.width = Pt(2)
        circle_positions.append((name, px, py, color, diam))

    # LAYER 4: ALL labels on top
    for name, px, py, color, diam in circle_positions:
        fs = 11 if 'TASUKI' in name else 9
        bld = 'TASUKI' in name
        if 'TASUKI' in name:
            label_y = py - diam // 2 - Inches(0.35)
        else:
            label_y = py + diam // 2 + Inches(0.05)
        add_text_box(slide, px - Inches(0.8), label_y, Inches(1.6), Inches(0.4),
                     name, font_size=fs, bold=bld, color=color, alignment=PP_ALIGN.CENTER)

    # LAYER 5: Title and axis labels
    add_text_box(slide, Inches(0.5), Inches(0.1), Inches(9), Inches(0.45),
                 '図6. TASUKIの理論的位置づけ', font_size=18,
                 bold=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.3),
                 '一人一票からの乖離度  →',
                 font_size=11, italic=True, color=C_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(0.9), Inches(0.7), Inches(0.5),
                 '←\n回顧的', font_size=10, italic=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.0), Inches(5.7), Inches(0.7), Inches(0.5),
                 '展望的\n→', font_size=10, italic=True, color=C_DARK,
                 alignment=PP_ALIGN.CENTER)

    # LAYER 6: Caption
    add_text_box(slide, Inches(0.3), Inches(7.05), Inches(9.4), Inches(0.4),
                 '図6. 既存の選挙制度改革提案に対するTASUKIの位置づけ。'
                 '横軸：OPOVからの乖離度。縦軸：時間的指向（回顧的〜展望的）。',
                 font_size=9, italic=True, color=C_GRAY, alignment=PP_ALIGN.CENTER)

def create_pptx(lang='en'):
    """Create the full PPTX for the given language."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    if lang == 'en':
        # Slide 1: Fig 1 - editable
        build_slide1_en(prs)
        # Slide 2: Fig 2 - image
        build_image_slide(prs, f'{FIGS}/fig2_influence_functions.png',
                          'Figure 2. Influence Function Family \u03c9(S)',
                          'Figure 2. (a) Different influence function specifications mapping accountability score S '
                          'to trust coefficient \u03c4. (b) Parameter sensitivity for the concave specification '
                          'across different [\u03c4_min, \u03c4_max] ranges.')
        # Slide 3: Fig 3 - image
        build_image_slide(prs, f'{FIGS}/fig3_simulation_results.png',
                          'Figure 3. Baseline Simulation Results',
                          'Figure 3. (a) Mean accountability score trajectory. (b) Evolutionary dynamics of candidate types. '
                          '(c) Trust coefficient distribution shift. (d) Voter welfare comparison across systems.')
        # Slide 4: Fig 4 - image
        build_image_slide(prs, f'{FIGS}/fig4_adversarial.png',
                          'Figure 4. Adversarial Robustness Analysis',
                          'Figure 4. (a) Exploitability index by strategy type and influence function. '
                          '(b) Genetic algorithm convergence for adversarial strategy search.')
        # Slide 5: Fig 5 - image
        build_image_slide(prs, f'{FIGS}/fig5_sensitivity.png',
                          'Figure 5. Sensitivity Analysis',
                          'Figure 5. (a) Heatmap of equilibrium accountability score as a function of '
                          '\u03c4_min and \u03c4_max. (b) Scalability analysis across candidate pool sizes.')
        # Slide 6: Fig 6 - editable
        build_slide6_en(prs)

        out_path = '/home/ubuntu/TASUKI_Figures_English.pptx'
    else:
        # Slide 1: Fig 1 - editable
        build_slide1_ja(prs)
        # Slide 2: Fig 2 - image
        build_image_slide(prs, f'{FIGS}/fig2_influence_functions.png',
                          '\u56f32. \u5f71\u97ff\u95a2\u6570\u65cf \u03c9(S)',
                          '\u56f32. (a) \u8aac\u660e\u8cac\u4efb\u5f97\u70b9S\u304b\u3089\u4fe1\u4efb\u4fc2\u6570\u03c4\u3078\u306e\u7570\u306a\u308b\u5f71\u97ff\u95a2\u6570\u4ed5\u69d8\u3002'
                          '(b) \u51f9\u95a2\u6570\u4ed5\u69d8\u306b\u304a\u3051\u308b[\u03c4_min, \u03c4_max]\u7bc4\u56f2\u306b\u5bfe\u3059\u308b\u30d1\u30e9\u30e1\u30fc\u30bf\u611f\u5ea6\u3002')
        # Slide 3: Fig 3 - image
        build_image_slide(prs, f'{FIGS}/fig3_simulation_results.png',
                          '\u56f33. \u30d9\u30fc\u30b9\u30e9\u30a4\u30f3\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\u7d50\u679c',
                          '\u56f33. (a) \u5e73\u5747\u8aac\u660e\u8cac\u4efb\u5f97\u70b9\u306e\u63a8\u79fb\u3002(b) \u5019\u88dc\u8005\u30bf\u30a4\u30d7\u306e\u9032\u5316\u52d5\u614b\u3002'
                          '(c) \u4fe1\u4efb\u4fc2\u6570\u5206\u5e03\u306e\u30b7\u30d5\u30c8\u3002(d) \u30b7\u30b9\u30c6\u30e0\u9593\u306e\u6709\u6a29\u8005\u539a\u751f\u6bd4\u8f03\u3002')
        # Slide 4: Fig 4 - image
        build_image_slide(prs, f'{FIGS}/fig4_adversarial.png',
                          '\u56f34. \u6575\u5bfe\u7684\u30ed\u30d0\u30b9\u30c8\u6027\u5206\u6790',
                          '\u56f34. (a) \u6226\u7565\u30bf\u30a4\u30d7\u3068\u5f71\u97ff\u95a2\u6570\u5225\u306e\u653b\u7565\u53ef\u80fd\u6027\u6307\u6a19\u3002'
                          '(b) \u6575\u5bfe\u7684\u6226\u7565\u63a2\u7d22\u306e\u907a\u4f1d\u7684\u30a2\u30eb\u30b4\u30ea\u30ba\u30e0\u53ce\u675f\u3002')
        # Slide 5: Fig 5 - image
        build_image_slide(prs, f'{FIGS}/fig5_sensitivity.png',
                          '\u56f35. \u611f\u5ea6\u5206\u6790',
                          '\u56f35. (a) \u03c4_min\u3068\u03c4_max\u306e\u95a2\u6570\u3068\u3057\u3066\u306e\u5747\u8861\u8aac\u660e\u8cac\u4efb\u5f97\u70b9\u306e\u30d2\u30fc\u30c8\u30de\u30c3\u30d7\u3002'
                          '(b) \u5019\u88dc\u8005\u30d7\u30fc\u30eb\u30b5\u30a4\u30ba\u5225\u306e\u30b9\u30b1\u30fc\u30e9\u30d3\u30ea\u30c6\u30a3\u5206\u6790\u3002')
        # Slide 6: Fig 6 - editable
        build_slide6_ja(prs)

        out_path = '/home/ubuntu/TASUKI_Figures_Japanese.pptx'

    prs.save(out_path)
    print(f'Saved: {out_path}')


if __name__ == '__main__':
    create_pptx('en')
    create_pptx('ja')
    print('All PPTX files created successfully!')
